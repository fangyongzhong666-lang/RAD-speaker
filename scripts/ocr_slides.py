#!/usr/bin/env python3
"""Extract and OCR images from PPTX slides when text extraction fails.

Usage:
  python ocr_slides.py input.pptx --output result.json
"""

import argparse
import json
import os
import sys
import tempfile
from pathlib import Path


def extract_slide_images(pptx_path: str) -> list[dict]:
    """Extract images from PPTX slides. Returns [{slide, image_path, text_found}].

    text_found=False means the slide is image-only and needs OCR.
    """
    try:
        from pptx import Presentation
    except ImportError:
        sys.stderr.write("python-pptx is required. pip install python-pptx\n")
        raise SystemExit(2)

    prs = Presentation(pptx_path)
    slides_data = []

    for i, slide in enumerate(prs.slides, 1):
        has_text = False
        images = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        has_text = True
                        break
            if shape.shape_type == 13:  # PICTURE
                img = shape.image
                ext = img.content_type.split("/")[-1] if "/" in img.content_type else "png"
                images.append({"blob": img.blob, "ext": ext})

        slides_data.append({
            "slide": i,
            "has_text": has_text,
            "image_count": len(images),
            "images": images,
        })

    return slides_data


def ocr_images(slides_data: list[dict], work_dir: str, languages: list[str]) -> list[dict]:
    """Run OCR on all images in slides that lack text. Returns OCR results per slide."""
    try:
        import easyocr
    except ImportError:
        sys.stderr.write("easyocr is required. pip install easyocr\n")
        raise SystemExit(2)

    # Set proxy from environment if available
    proxy = os.environ.get("HTTP_PROXY") or os.environ.get("HTTPS_PROXY")
    if proxy:
        os.environ.setdefault("HTTP_PROXY", proxy)
        os.environ.setdefault("HTTPS_PROXY", proxy)

    reader = easyocr.Reader(languages, gpu=False)
    results = []

    for slide_data in slides_data:
        slide_num = slide_data["slide"]
        slide_result = {"slide": slide_num, "ocr_texts": []}

        if slide_data["has_text"] and slide_data["image_count"] == 0:
            slide_result["status"] = "skipped (has text, no images)"
            results.append(slide_result)
            continue

        for idx, img in enumerate(slide_data["images"]):
            img_path = Path(work_dir) / f"slide_{slide_num:02d}_img_{idx}.{img['ext']}"
            img_path.write_bytes(img["blob"])

            try:
                texts = reader.readtext(str(img_path), detail=0)
                for line in texts:
                    slide_result["ocr_texts"].append(line)
            except Exception as e:
                slide_result.setdefault("errors", []).append(
                    f"slide {slide_num} img {idx}: {e}"
                )

        slide_result["status"] = "ocr complete" if slide_result["ocr_texts"] else "no text found"
        results.append(slide_result)

    return results


def main(argv: list[str]) -> int:
    parser = argparse.ArgumentParser(
        description="OCR images from PPTX slides and return extracted text."
    )
    parser.add_argument("input", type=Path, help="Input .pptx file")
    parser.add_argument("--output", "-o", type=Path, help="Output JSON file with OCR results")
    parser.add_argument("--languages", "-l", default="ch_sim,en",
                        help="Comma-separated easyocr language codes (default: ch_sim,en)")
    parser.add_argument("--force", action="store_true",
                        help="OCR all slides even if text is present")
    args = parser.parse_args(argv)

    if not args.input.exists():
        sys.stderr.write(f"Input not found: {args.input}\n")
        return 1

    languages = [lang.strip() for lang in args.languages.split(",")]

    # Extract slide images
    print(f"Extracting images from {args.input}...", file=sys.stderr)
    slides_data = extract_slide_images(str(args.input))

    total_images = sum(s["image_count"] for s in slides_data)
    textless_slides = sum(1 for s in slides_data if not s["has_text"])
    print(f"Found {len(slides_data)} slides, {total_images} images, "
          f"{textless_slides} slides without text", file=sys.stderr)

    # Filter to slides needing OCR
    if not args.force:
        slides_to_ocr = [s for s in slides_data if not s["has_text"] or s["image_count"] > 0]
    else:
        slides_to_ocr = slides_data

    # Run OCR
    if slides_to_ocr:
        work_dir = args.output.parent / "ocr_work" if args.output else Path(tempfile.mkdtemp())
        work_dir.mkdir(parents=True, exist_ok=True)
        print(f"Running OCR on {len(slides_to_ocr)} slides...", file=sys.stderr)
        ocr_results = ocr_images(slides_to_ocr, str(work_dir), languages)
    else:
        ocr_results = []

    # Build output
    output = {
        "file": str(args.input),
        "total_slides": len(slides_data),
        "slides_without_text": textless_slides,
        "total_images_ocr": total_images,
        "ocr_results": ocr_results,
    }

    if args.output:
        args.output.write_text(json.dumps(output, ensure_ascii=False, indent=2), encoding="utf-8")
        print(f"Saved: {args.output}", file=sys.stderr)
    else:
        print(json.dumps(output, ensure_ascii=False, indent=2))

    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
